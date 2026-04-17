from __future__ import annotations

import io
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

from .export_context import build_export_context
from .views import _chart_to_png

W = Inches(13.333)
H = Inches(7.5)
DARK = RGBColor(11, 9, 24)
SURF = RGBColor(17, 24, 39)
VIOLET = RGBColor(124, 58, 237)
WHITE = RGBColor(255,255,255)
MUTED = RGBColor(148,163,184)
GREEN = RGBColor(16,185,129)
AMBER = RGBColor(245,158,11)
ROSE = RGBColor(244,63,94)


def _rgb(rgb_tuple):
    return RGBColor(*rgb_tuple)


def _bg(slide, color=DARK):
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background();
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def _box(slide, l, t, w, h, text='', size=14, bold=False, color=WHITE, fill=None, align=PP_ALIGN.LEFT):
    if fill is not None:
        shp = slide.shapes.add_shape(1, l, t, w, h)
        shp.fill.solid(); shp.fill.fore_color.rgb = fill; shp.line.fill.background()
    tx = slide.shapes.add_textbox(l, t, w, h)
    p = tx.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run(); run.text = str(text); run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = color; run.font.name = 'Calibri'
    return tx


def build_pptx(upload, request=None) -> bytes:
    ctx = build_export_context(upload, mode='executive', request=request)
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H; blank = prs.slide_layouts[6]
    hero = (ctx.get('dashboard') or {}).get('hero') or {}
    polish = ctx.get('presentation_polish') or {}
    branding = ctx.get('branding') or {}
    dark = _rgb(branding.get('bg_rgb') or (11,9,24))
    surf = _rgb(branding.get('surface_rgb') or (17,24,39))
    violet = _rgb(branding.get('primary_rgb') or (124,58,237))
    accent = _rgb(branding.get('accent_rgb') or (59,130,246))
    # cover
    s = prs.slides.add_slide(blank); _bg(s, dark)
    logo_path = branding.get('logo_path') or ''
    if logo_path:
        try:
            s.shapes.add_picture(logo_path, Inches(0.45), Inches(0.45), height=Inches(0.7))
        except Exception:
            pass
    _box(s, Inches(0.45), Inches(1.1), Inches(12), Inches(0.45), branding.get('report_title') or upload.original_name, 24, True, color=WHITE)
    _box(s, Inches(0.45), Inches(1.65), Inches(12), Inches(0.55), upload.original_name, 18, True, color=accent)
    _box(s, Inches(0.45), Inches(2.25), Inches(10), Inches(0.5), hero.get('headline') or 'Board export', 16, False, MUTED)
    _box(s, Inches(0.45), Inches(2.85), Inches(11.5), Inches(1.1), hero.get('subheadline') or 'Decision-ready export generated from the current dashboard state.', 14)
    notes = polish.get('export_notes') or []
    for idx, note in enumerate(notes[:3]):
        _box(s, Inches(0.55), Inches(4.15 + idx*0.34), Inches(11.6), Inches(0.28), f'• {note}', 11, False, MUTED)
    _box(s, Inches(0.45), Inches(6.85), Inches(12), Inches(0.25), branding.get('footer_text') or '', 10, False, MUTED)
    _box(s, Inches(10.8), Inches(0.35), Inches(2.0), Inches(0.25), branding.get('watermark') or '', 10, True, MUTED, align=PP_ALIGN.RIGHT)

    # executive summary
    s = prs.slides.add_slide(blank); _bg(s, dark)
    _box(s, Inches(0.4), Inches(0.25), Inches(12), Inches(0.5), 'Executive summary', 22, True)
    cards = ctx.get('decision_cards') or []
    fills = [ROSE, AMBER, GREEN, VIOLET]
    for i, card in enumerate(cards[:4]):
        left = Inches(0.45 + (i%2)*6.25)
        top = Inches(1.0 + (i//2)*2.6)
        fill = fills[i%len(fills)] if i else accent
        _box(s, left, top, Inches(5.8), Inches(2.2), '', fill=surf)
        _box(s, left+Inches(0.15), top+Inches(0.15), Inches(5.4), Inches(0.3), f"{card.get('emoji','')} {card.get('label','Signal')}", 10, True, fill)
        _box(s, left+Inches(0.15), top+Inches(0.55), Inches(5.4), Inches(0.45), card.get('title',''), 16, True)
        _box(s, left+Inches(0.15), top+Inches(1.0), Inches(5.4), Inches(0.65), card.get('body',''), 11)
        _box(s, left+Inches(0.15), top+Inches(1.65), Inches(5.4), Inches(0.35), f"Action: {card.get('action','')}", 10, False, MUTED)

    # sections
    for section in (ctx.get('sections') or [])[:5]:
        s = prs.slides.add_slide(blank); _bg(s, dark)
        _box(s, Inches(0.4), Inches(0.25), Inches(12), Inches(0.5), section.get('title') or 'Section', 22, True, color=accent)
        style = (polish.get('section_styles') or {}).get(section.get('key') or '', {})
        _box(s, Inches(0.4), Inches(0.8), Inches(12), Inches(0.45), section.get('intro') or '', 12, False, MUTED)
        charts = section.get('charts') or []
        positions = [(Inches(0.45), Inches(1.35)), (Inches(6.75), Inches(1.35))]
        for (left, top), chart in zip(positions, charts[:2]):
            png = _chart_to_png(chart, width_in=5.6, height_in=2.9)
            if png:
                s.shapes.add_picture(io.BytesIO(png), left, top, width=Inches(5.9), height=Inches(3.15))
            meta = (polish.get('chart_styles') or {}).get(str(getattr(chart, 'id', '')), {})
            prefix = '★ ' if meta.get('priority') == 'primary' else ('◆ ' if meta.get('priority') == 'secondary' else '')
            _box(s, left, top+Inches(3.18), Inches(5.9), Inches(0.35), f'{prefix}{chart.title}', 11, True)
            explanation = getattr(chart, 'explanation', None) or {}
            _box(s, left, top+Inches(3.55), Inches(5.9), Inches(1.1), explanation.get('summary') or explanation.get('why_it_matters') or '', 10)

    # exception tables
    for table in (ctx.get('exception_tables') or [])[:3]:
        s = prs.slides.add_slide(blank); _bg(s, dark)
        _box(s, Inches(0.4), Inches(0.25), Inches(12), Inches(0.5), table.get('title') or 'Decision table', 22, True)
        _box(s, Inches(0.4), Inches(0.8), Inches(12), Inches(0.35), table.get('intro') or '', 11, False, MUTED)
        cols = table.get('columns') or ['Entity','Metric','Support']
        rows = table.get('rows') or []
        headers = '   |   '.join(cols)
        _box(s, Inches(0.45), Inches(1.35), Inches(12), Inches(0.3), headers, 11, True, VIOLET)
        for i, row in enumerate(rows[:8]):
            y = Inches(1.75 + i*0.55)
            _box(s, Inches(0.45), y, Inches(3.0), Inches(0.3), row.get('entity',''), 11)
            _box(s, Inches(3.6), y, Inches(2.0), Inches(0.3), row.get('metric',''), 11)
            _box(s, Inches(5.8), y, Inches(6.8), Inches(0.3), row.get('support',''), 10, False, MUTED)

    # appendix
    s = prs.slides.add_slide(blank); _bg(s, dark)
    profile = ctx.get('profile') or {}
    governance = ctx.get('governance') or {}
    _box(s, Inches(0.4), Inches(0.25), Inches(12), Inches(0.5), 'Appendix, audit & methodology', 22, True)
    lines = [
        f"Source file: {upload.original_name}",
        f"Rows: {profile.get('row_count') or upload.row_count}",
        f"Columns: {profile.get('column_count') or upload.column_count}",
        f"Analysis type: {((profile.get('analysis_classification') or {}).get('analysis_type') or '').replace('_',' ').title()}",
        f"Primary dimension: {((profile.get('analysis_classification') or {}).get('primary_dimension') or '')}",
        f"Primary measure: {((profile.get('analysis_classification') or {}).get('primary_measure') or '')}",
    ]
    for i, line in enumerate(lines):
        _box(s, Inches(0.6), Inches(1.0 + i*0.42), Inches(10.5), Inches(0.28), line, 13)
    start = 1.0 + len(lines)*0.42 + 0.25
    _box(s, Inches(0.6), Inches(start), Inches(4.0), Inches(0.25), 'Methodology', 12, True, VIOLET)
    for idx, item in enumerate((governance.get('methodology') or [])[:3]):
        _box(s, Inches(0.8), Inches(start + 0.28 + idx*0.3), Inches(5.6), Inches(0.24), f'• {item}', 10)
    _box(s, Inches(6.5), Inches(start), Inches(4.0), Inches(0.25), 'Known caveats', 12, True, AMBER)
    caveats = governance.get('caveats') or ['No major caveats flagged.']
    for idx, item in enumerate(caveats[:4]):
        _box(s, Inches(6.7), Inches(start + 0.28 + idx*0.3), Inches(5.8), Inches(0.24), f'• {item}', 10)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0); return buf.getvalue()
