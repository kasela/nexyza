"""
Background export tasks — run by Django-Q2 worker.
"""
import io
import os
import logging
from django.conf import settings
from django.utils import timezone

logger = logging.getLogger(__name__)


def run_export(job_id: str):
    """Execute an export job. Called by Django-Q2."""
    from .models import ExportJob

    try:
        job = ExportJob.objects.select_related('upload', 'user').get(pk=job_id)
    except ExportJob.DoesNotExist:
        return

    try:
        # Get branding for the user
        branding = None
        try:
            branding = job.user.branding
        except Exception:
            pass

        if job.fmt == 'pdf':
            data = _build_pdf(job.upload, branding)
            ext  = 'pdf'
            mime = 'application/pdf'
        elif job.fmt == 'pptx':
            data = _build_pptx(job.upload, branding, theme=job.theme)
            ext  = 'pptx'
            mime = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
        else:
            raise ValueError(f'Unknown format: {job.fmt}')

        # Save to media
        fname    = f'exports/{job.user_id}/{job_id}.{ext}'
        full_dir = os.path.join(settings.MEDIA_ROOT, 'exports', str(job.user_id))
        os.makedirs(full_dir, exist_ok=True)
        with open(os.path.join(settings.MEDIA_ROOT, fname), 'wb') as f:
            f.write(data)

        job.result_url = f'{settings.MEDIA_URL}{fname}'
        job.status     = ExportJob.STATUS_DONE
        job.save(update_fields=['result_url', 'status', 'updated_at'])

    except Exception as e:
        logger.error(f'Export job {job_id} failed: {e}', exc_info=True)
        job.status = ExportJob.STATUS_ERROR
        job.error  = str(e)[:500]
        job.save(update_fields=['status', 'error', 'updated_at'])


# ── PDF builder ───────────────────────────────────────────────────────────────

def _hex_to_rgb(hex_color: str):
    """Convert '#7c3aed' → (0x7c, 0x3a, 0xed)."""
    h = hex_color.lstrip('#')
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))


def _build_pdf(upload, branding=None) -> bytes:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.lib import colors
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                     Table, TableStyle, HRFlowable, Image,
                                     KeepTogether, PageBreak)
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT
    from reportlab.pdfgen import canvas as rl_canvas
    from apps.analyser.models import ChartConfig
    from apps.exports.views import _chart_to_png

    # ── Branding colours ──────────────────────────────────────────────────────
    primary_hex = (branding.primary_color if branding else None) or '#7c3aed'
    accent_hex  = (branding.accent_color  if branding else None) or '#3b82f6'
    app_name    = (branding.app_name      if branding else None) or 'Nexyza'

    PRIMARY = colors.HexColor(primary_hex)
    ACCENT  = colors.HexColor(accent_hex)
    DARK    = colors.HexColor('#0f172a')
    LIGHT   = colors.HexColor('#f8fafc')
    MID     = colors.HexColor('#64748b')
    WHITE   = colors.white
    BLACK   = colors.HexColor('#1e293b')

    analysis = upload.analysis_result or {}
    cols_data = analysis.get('columns', [])

    buf = io.BytesIO()

    # ── Page numbering + header/footer ────────────────────────────────────────
    page_num = [0]

    def on_page(canv, doc):
        page_num[0] += 1
        w, h = A4
        canv.saveState()
        # Footer bar
        canv.setFillColor(colors.HexColor('#f1f5f9'))
        canv.rect(0, 0, w, 1.2*cm, fill=1, stroke=0)
        canv.setFillColor(PRIMARY)
        canv.rect(0, 0, w, 0.25*cm, fill=1, stroke=0)
        # Page number
        canv.setFont('Helvetica', 8)
        canv.setFillColor(MID)
        canv.drawCentredString(w/2, 0.45*cm, f'Page {page_num[0]}')
        canv.drawString(2*cm, 0.45*cm, app_name)
        canv.drawRightString(w - 2*cm, 0.45*cm, upload.original_name[:50])
        # Top accent line
        canv.setFillColor(PRIMARY)
        canv.rect(0, h - 0.25*cm, w, 0.25*cm, fill=1, stroke=0)
        canv.restoreState()

    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        leftMargin=2.2*cm, rightMargin=2.2*cm,
        topMargin=2.5*cm, bottomMargin=2*cm,
        onFirstPage=on_page, onLaterPages=on_page,
    )

    # ── Styles ────────────────────────────────────────────────────────────────
    styles = getSampleStyleSheet()
    title_s = ParagraphStyle('T',  parent=styles['Title'],
                              textColor=PRIMARY, fontSize=22, spaceAfter=8, spaceBefore=8)
    sub_s   = ParagraphStyle('S',  parent=styles['Normal'],
                              textColor=MID, fontSize=10, spaceAfter=4)
    intro_s = ParagraphStyle('I',  parent=styles['Normal'],
                              textColor=BLACK, fontSize=10.5, spaceAfter=8, leading=16)
    h2_s    = ParagraphStyle('H2', parent=styles['Heading2'],
                              textColor=PRIMARY, fontSize=14,
                              spaceBefore=18, spaceAfter=6, borderPad=0)
    h3_s    = ParagraphStyle('H3', parent=styles['Heading3'],
                              textColor=colors.HexColor(accent_hex),
                              fontSize=11, spaceBefore=10, spaceAfter=4)
    body_s  = ParagraphStyle('B',  parent=styles['Normal'],
                              textColor=BLACK, fontSize=9.5, spaceAfter=4, leading=14)
    caption_s = ParagraphStyle('C', parent=styles['Normal'],
                                textColor=MID, fontSize=8, spaceAfter=6,
                                alignment=TA_CENTER)
    toc_s   = ParagraphStyle('TOC', parent=styles['Normal'],
                              textColor=BLACK, fontSize=10, spaceAfter=4, leading=16)

    story = []

    # ── Cover page ─────────────────────────────────────────────────────────────
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph(upload.original_name, title_s))
    story.append(Paragraph('Data Analysis Report', ParagraphStyle(
        'SUBT', parent=styles['Normal'], textColor=ACCENT, fontSize=13,
        spaceAfter=4)))
    story.append(HRFlowable(color=PRIMARY, thickness=2, spaceAfter=14, width='100%'))
    story.append(Paragraph(
        f"<b>Dataset:</b> {upload.original_name}<br/>"
        f"<b>Rows:</b> {upload.row_count:,} &nbsp;&nbsp; "
        f"<b>Columns:</b> {upload.column_count} &nbsp;&nbsp; "
        f"<b>Format:</b> {upload.file_type.upper()}<br/>"
        f"<b>Generated:</b> {timezone.now().strftime('%B %d, %Y')} &nbsp;&nbsp; "
        f"<b>Prepared by:</b> {app_name}",
        body_s))
    story.append(PageBreak())

    # ── Table of Contents ─────────────────────────────────────────────────────
    story.append(Paragraph('Table of Contents', h2_s))
    story.append(HRFlowable(color=PRIMARY, thickness=1, spaceAfter=10, width='100%'))
    toc_items = [
        ('1.', 'Introduction & Dataset Overview'),
        ('2.', 'Column Statistics'),
        ('3.', 'Data Visualisations'),
        ('4.', 'AI-Generated Insights'),
    ]
    for num, title in toc_items:
        story.append(Paragraph(
            f'<font color="{primary_hex}">{num}</font>&nbsp;&nbsp;{title}',
            toc_s))
    story.append(PageBreak())

    # ── 1. Introduction ────────────────────────────────────────────────────────
    story.append(Paragraph('1. Introduction & Dataset Overview', h2_s))
    story.append(HRFlowable(color=PRIMARY, thickness=0.5, spaceAfter=8, width='100%'))
    numeric_cols = [c for c in cols_data if c.get('is_numeric')]
    text_cols    = [c for c in cols_data if not c.get('is_numeric')]
    null_pcts    = [c.get('null_pct', 0) for c in cols_data]
    avg_null     = sum(null_pcts) / len(null_pcts) if null_pcts else 0

    story.append(Paragraph(
        f"This report presents an automated analysis of <b>{upload.original_name}</b>, "
        f"a {upload.file_type.upper()} dataset containing <b>{upload.row_count:,} rows</b> "
        f"and <b>{upload.column_count} columns</b>. "
        f"The dataset includes {len(numeric_cols)} numeric column{'s' if len(numeric_cols) != 1 else ''} "
        f"and {len(text_cols)} categorical column{'s' if len(text_cols) != 1 else ''}. "
        f"Average null rate across all columns is {avg_null:.1f}%.",
        intro_s))

    # Summary KPI table
    kpi_data = [
        ['Metric', 'Value'],
        ['Total Rows',           f"{upload.row_count:,}"],
        ['Total Columns',        str(upload.column_count)],
        ['Numeric Columns',      str(len(numeric_cols))],
        ['Categorical Columns',  str(len(text_cols))],
        ['Average Null Rate',    f"{avg_null:.1f}%"],
        ['File Format',          upload.file_type.upper()],
    ]
    kpi_tbl = Table(kpi_data, colWidths=[8*cm, 7*cm])
    kpi_tbl.setStyle(TableStyle([
        ('BACKGROUND',  (0,0), (-1,0), PRIMARY),
        ('TEXTCOLOR',   (0,0), (-1,0), WHITE),
        ('FONTNAME',    (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE',    (0,0), (-1,-1), 9.5),
        ('FONTNAME',    (0,1), (-1,-1), 'Helvetica'),
        ('TEXTCOLOR',   (0,1), (-1,-1), BLACK),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [LIGHT, WHITE]),
        ('GRID',        (0,0), (-1,-1), 0.3, colors.HexColor('#cbd5e1')),
        ('TOPPADDING',  (0,0), (-1,-1), 5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ('LEFTPADDING', (0,0), (-1,-1), 8),
    ]))
    story.append(kpi_tbl)
    story.append(PageBreak())

    # ── 2. Column Statistics ───────────────────────────────────────────────────
    story.append(Paragraph('2. Column Statistics', h2_s))
    story.append(HRFlowable(color=PRIMARY, thickness=0.5, spaceAfter=8, width='100%'))
    tdata = [['Column', 'Type', 'Nulls %', 'Unique', 'Min', 'Max', 'Mean']]
    for col in cols_data:
        np_ = col.get('null_pct', 0)
        tdata.append([
            col['name'][:28],
            'Numeric' if col.get('is_numeric') else 'Text',
            f"{np_:.1f}%",
            f"{col.get('unique_count', ''):,}" if col.get('unique_count') else '—',
            str(col.get('min', '—'))[:12],
            str(col.get('max', '—'))[:12],
            f"{col.get('mean', 0):.3f}" if col.get('is_numeric') and col.get('mean') is not None else '—',
        ])

    col_widths = [5*cm, 1.8*cm, 1.8*cm, 1.6*cm, 2.1*cm, 2.1*cm, 2.1*cm]
    tbl = Table(tdata, colWidths=col_widths)

    row_styles = [TableStyle([
        ('BACKGROUND',  (0,0), (-1,0), PRIMARY),
        ('TEXTCOLOR',   (0,0), (-1,0), WHITE),
        ('FONTNAME',    (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTNAME',    (0,1), (-1,-1), 'Helvetica'),
        ('FONTSIZE',    (0,0), (-1,-1), 8.5),
        ('TEXTCOLOR',   (0,1), (-1,-1), BLACK),
        ('ROWBACKGROUNDS', (0,1), (-1,-1), [LIGHT, WHITE]),
        ('GRID',        (0,0), (-1,-1), 0.25, colors.HexColor('#e2e8f0')),
        ('ALIGN',       (1,0), (-1,-1), 'CENTER'),
        ('ALIGN',       (0,0), (0,-1), 'LEFT'),
        ('TOPPADDING',  (0,0), (-1,-1), 4),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
    ])]

    # Highlight high null cells in red
    for ri, col in enumerate(cols_data, 1):
        if col.get('null_pct', 0) > 20:
            row_styles.append(TableStyle([
                ('TEXTCOLOR', (2, ri), (2, ri), colors.HexColor('#dc2626')),
                ('FONTNAME',  (2, ri), (2, ri), 'Helvetica-Bold'),
            ]))

    for s in row_styles:
        tbl.setStyle(s)

    story.append(tbl)
    story.append(PageBreak())

    # ── 3. Charts ──────────────────────────────────────────────────────────────
    story.append(Paragraph('3. Data Visualisations', h2_s))
    story.append(HRFlowable(color=PRIMARY, thickness=0.5, spaceAfter=8, width='100%'))

    charts = list(ChartConfig.objects.filter(upload=upload).order_by('sort_order', 'created_at'))
    renderable = [c for c in charts if c.cached_data and (
        c.cached_data.get('labels') or c.chart_type == 'kpi' or c.cached_data.get('kpi'))]

    page_w  = A4[0] - 4.4*cm
    chart_h = page_w * 0.42

    if not renderable:
        story.append(Paragraph('<i>No charts available for this dataset.</i>', body_s))
    else:
        for i, chart in enumerate(renderable):
            story.append(Paragraph(chart.title, h3_s))
            insight = (chart.config_json or {}).get('insight', '')
            if insight:
                story.append(Paragraph(insight, caption_s))
            png = _chart_to_png(chart, branding=branding,
                                width_in=page_w/28.35, height_in=chart_h/28.35)
            if png:
                img_buf = io.BytesIO(png)
                story.append(Image(img_buf, width=page_w, height=chart_h))
            story.append(Spacer(1, 10))
            if (i+1) % 2 == 0 and i < len(renderable)-1:
                story.append(PageBreak())

    story.append(PageBreak())

    # ── 4. AI Insights ─────────────────────────────────────────────────────────
    story.append(Paragraph('4. AI-Generated Insights', h2_s))
    story.append(HRFlowable(color=PRIMARY, thickness=0.5, spaceAfter=8, width='100%'))
    if upload.ai_insights:
        for line in upload.ai_insights.split('\n'):
            line = line.strip().lstrip('#*-').strip()
            if line:
                style = h3_s if line.startswith('##') or len(line) < 60 else body_s
                story.append(Paragraph(line.lstrip('#').strip(), style))
    else:
        story.append(Paragraph(
            '<i>AI insights are available on Plus and Pro plans. '
            'Upgrade at nexyza.com/billing/pricing/ to unlock AI-generated analysis.</i>',
            body_s))

    story.append(PageBreak())

    # ── Thank you / closing page ───────────────────────────────────────────────
    story.append(Spacer(1, 4*cm))
    story.append(Paragraph('Thank You', ParagraphStyle(
        'TY', parent=styles['Title'], textColor=PRIMARY, fontSize=26,
        alignment=TA_CENTER, spaceAfter=12)))
    story.append(Paragraph(
        'This report was generated automatically by Nexyza.',
        ParagraphStyle('TC', parent=styles['Normal'], textColor=MID,
                       fontSize=11, alignment=TA_CENTER, spaceAfter=6)))
    story.append(Paragraph(
        f'<a href="https://nexyza.com" color="{primary_hex}">nexyza.com</a>',
        ParagraphStyle('TL', parent=styles['Normal'], textColor=ACCENT,
                       fontSize=11, alignment=TA_CENTER)))

    doc.build(story)
    buf.seek(0)
    return buf.getvalue()


# ── PPTX builder ──────────────────────────────────────────────────────────────

PPTX_THEMES = {
    'dark': {
        'bg':      '0F0C1E',
        'surface': '1E1B2E',
        'text':    'E2E0F0',
        'muted':   '94A3B8',
        'border':  '334155',
    },
    'light': {
        'bg':      'F8FAFC',
        'surface': 'FFFFFF',
        'text':    '1E293B',
        'muted':   '64748B',
        'border':  'E2E8F0',
    },
    'corporate': {
        'bg':      '1E3A5F',
        'surface': '1A3354',
        'text':    'F0F4F8',
        'muted':   '94A3B8',
        'border':  '2D4A6B',
    },
    'minimal': {
        'bg':      'FFFFFF',
        'surface': 'F1F5F9',
        'text':    '0F172A',
        'muted':   '475569',
        'border':  'CBD5E1',
    },
}


def _rgb(hex6: str):
    from pptx.dml.color import RGBColor
    h = hex6.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def _build_pptx(upload, branding=None, theme: str = 'dark') -> bytes:
    import io as _io
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from apps.analyser.models import ChartConfig
    from apps.exports.pptx_export import _chart_png  # reuse existing renderer

    # Branding
    primary_hex = (branding.primary_color if branding else None) or '#7c3aed'
    accent_hex  = (branding.accent_color  if branding else None) or '#3b82f6'
    app_name    = (branding.app_name      if branding else None) or 'Nexyza'

    T = PPTX_THEMES.get(theme, PPTX_THEMES['dark'])

    C_BG      = _rgb(T['bg'])
    C_SURFACE = _rgb(T['surface'])
    C_TEXT    = _rgb(T['text'])
    C_MUTED   = _rgb(T['muted'])
    C_PRIMARY = _rgb(primary_hex)
    C_ACCENT  = _rgb(accent_hex)

    analysis  = upload.analysis_result or {}
    cols_data = analysis.get('columns', [])
    numeric   = [c for c in cols_data if c.get('is_numeric')]

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    W, H = Inches(13.33), Inches(7.5)
    blank = prs.slide_layouts[6]

    def add_bg(slide, color=None):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color or C_BG

    def add_rect(slide, l, t, w, h, color):
        from pptx.util import Pt as _Pt
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid(); s.fill.fore_color.rgb = color
        s.line.fill.background()
        return s

    def add_txt(slide, l, t, w, h, text, size, bold=False,
                color=None, align=PP_ALIGN.LEFT, wrap=True):
        color = color or C_TEXT
        txb = slide.shapes.add_textbox(l, t, w, h)
        tf  = txb.text_frame; tf.word_wrap = wrap
        p   = tf.paragraphs[0]; p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = 'Calibri'
        return txb

    charts = list(ChartConfig.objects.filter(upload=upload).order_by('sort_order','created_at'))
    renderable = [c for c in charts if c.cached_data and (
        c.cached_data.get('labels') or c.chart_type == 'kpi' or c.cached_data.get('kpi'))]

    # ── Slide 1: Title / Welcome ───────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    add_bg(s1)
    add_rect(s1, 0, 0, W, Inches(0.12), C_PRIMARY)
    add_rect(s1, 0, H - Inches(0.12), W, Inches(0.12), C_PRIMARY)
    add_rect(s1, 0, 0, Inches(0.12), H, C_PRIMARY)

    add_txt(s1, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
            app_name.upper(), 11, bold=True, color=C_PRIMARY, align=PP_ALIGN.LEFT)
    add_txt(s1, Inches(1), Inches(1.9), Inches(11), Inches(1.6),
            upload.original_name, 36, bold=True, color=C_TEXT)
    add_txt(s1, Inches(1), Inches(3.6), Inches(11), Inches(0.5),
            'Data Analysis Report', 18, color=C_ACCENT)
    add_txt(s1, Inches(1), Inches(4.3), Inches(11), Inches(0.4),
            f"Generated {timezone.now().strftime('%B %d, %Y')} · {upload.row_count:,} rows · {upload.column_count} columns",
            12, color=C_MUTED)

    # ── Slide 2: Agenda / Contents ────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    add_bg(s2)
    add_rect(s2, 0, 0, W, Inches(0.12), C_PRIMARY)
    add_txt(s2, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
            'Agenda', 24, bold=True, color=C_PRIMARY)
    items = [
        ('01', 'Introduction & Dataset Overview'),
        ('02', 'Column Statistics & Data Quality'),
        ('03', 'Charts & Visualisations'),
        ('04', 'AI-Generated Insights'),
        ('05', 'Q&A'),
    ]
    for i, (num, label) in enumerate(items):
        y = Inches(1.3 + i * 1.05)
        add_rect(s2, Inches(0.8), y + Inches(0.1), Inches(0.6), Inches(0.6), C_PRIMARY)
        add_txt(s2, Inches(0.8), y + Inches(0.08), Inches(0.6), Inches(0.6),
                num, 16, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_txt(s2, Inches(1.6), y + Inches(0.08), Inches(10), Inches(0.6),
                label, 16, color=C_TEXT)

    # ── Slide 3: Introduction ─────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    add_bg(s3)
    add_rect(s3, 0, 0, W, Inches(0.12), C_PRIMARY)
    add_txt(s3, Inches(0.8), Inches(0.25), Inches(12), Inches(0.55),
            'Introduction', 22, bold=True, color=C_PRIMARY)
    add_txt(s3, Inches(0.8), Inches(1.0), Inches(11.5), Inches(1.5),
            f"This presentation covers an automated analysis of {upload.original_name}. "
            f"The dataset contains {upload.row_count:,} rows across {upload.column_count} columns — "
            f"{len(numeric)} numeric and {len(cols_data)-len(numeric)} categorical.",
            13, color=C_TEXT, wrap=True)

    kpis = [
        (f"{upload.row_count:,}", "Rows"),
        (str(upload.column_count), "Columns"),
        (str(len(numeric)), "Numeric Cols"),
        (upload.file_type.upper(), "Format"),
    ]
    bw = Inches(2.8)
    for i, (val, lbl) in enumerate(kpis):
        lx = Inches(0.6 + i * 3.1)
        add_rect(s3, lx, Inches(2.8), bw, Inches(2.2), C_SURFACE)
        add_rect(s3, lx, Inches(2.8), bw, Inches(0.08), C_PRIMARY)
        add_txt(s3, lx, Inches(3.0), bw, Inches(1.0), val,
                28, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
        add_txt(s3, lx, Inches(4.0), bw, Inches(0.4), lbl,
                11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # ── Slide 4: Column stats table ───────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    add_bg(s4)
    add_rect(s4, 0, 0, W, Inches(0.12), C_PRIMARY)
    add_txt(s4, Inches(0.8), Inches(0.25), Inches(12), Inches(0.55),
            'Column Statistics', 22, bold=True, color=C_PRIMARY)

    hdrs  = ['Column', 'Type', 'Nulls%', 'Unique', 'Min', 'Max', 'Mean']
    cws   = [Inches(3.2), Inches(1.1), Inches(1.1), Inches(1.0),
             Inches(1.7), Inches(1.7), Inches(1.7)]
    rh, y0, x0 = Inches(0.42), Inches(1.0), Inches(0.5)

    x = x0
    for hd, cw in zip(hdrs, cws):
        add_rect(s4, x, y0, cw - Inches(0.04), rh, C_PRIMARY)
        add_txt(s4, x + Inches(0.06), y0 + Inches(0.06), cw - Inches(0.12),
                rh - Inches(0.1), hd, 9, bold=True, color=C_TEXT)
        x += cw

    for ri, col in enumerate(cols_data[:12]):
        rc = C_SURFACE if ri % 2 == 0 else C_BG
        y  = y0 + rh * (ri + 1)
        x  = x0
        np_ = col.get('null_pct', 0)
        row = [
            col['name'][:30],
            'num' if col.get('is_numeric') else 'text',
            f"{np_:.1f}%",
            str(col.get('unique_count', '—')),
            str(col.get('min', '—'))[:12],
            str(col.get('max', '—'))[:12],
            f"{col.get('mean',0):.2f}" if col.get('is_numeric') and col.get('mean') is not None else '—',
        ]
        for vi, (val, cw) in enumerate(zip(row, cws)):
            add_rect(s4, x, y, cw - Inches(0.04), rh, rc)
            tc = _rgb('#dc2626') if vi == 2 and np_ > 20 else C_MUTED if vi > 0 else C_TEXT
            add_txt(s4, x + Inches(0.06), y + Inches(0.06),
                    cw - Inches(0.12), rh - Inches(0.1), val, 8, color=tc)
            x += cw

    # ── Chart slides ──────────────────────────────────────────────────────────
    for chart in renderable:
        sc = prs.slides.add_slide(blank)
        add_bg(sc)
        add_rect(sc, 0, 0, W, Inches(0.12), C_PRIMARY)
        add_txt(sc, Inches(0.6), Inches(0.2), Inches(12.1), Inches(0.55),
                chart.title, 20, bold=True, color=C_TEXT)
        insight = (chart.config_json or {}).get('insight', '')
        if insight:
            add_txt(sc, Inches(0.6), Inches(0.82), Inches(12.1), Inches(0.35),
                    insight, 10, color=C_MUTED)

        png_buf = _chart_png(chart, w_in=12.0, h_in=5.4,
                              primary=primary_hex, accent=accent_hex,
                              theme=T)
        if png_buf:
            top = Inches(1.25) if insight else Inches(0.95)
            sc.shapes.add_picture(png_buf, Inches(0.6), top,
                                   width=Inches(12.0), height=H - top - Inches(0.3))

    # ── AI Insights slide ─────────────────────────────────────────────────────
    if upload.ai_insights:
        si = prs.slides.add_slide(blank)
        add_bg(si)
        add_rect(si, 0, 0, W, Inches(0.12), C_PRIMARY)
        add_txt(si, Inches(0.8), Inches(0.25), Inches(12), Inches(0.55),
                'AI-Generated Insights', 22, bold=True, color=C_PRIMARY)
        add_txt(si, Inches(0.8), Inches(1.0), Inches(11.5), Inches(5.8),
                upload.ai_insights[:1800], 11, color=C_TEXT, wrap=True)

    # ── Q&A slide ──────────────────────────────────────────────────────────────
    sq = prs.slides.add_slide(blank)
    add_bg(sq)
    add_rect(sq, 0, 0, W, Inches(0.12), C_PRIMARY)
    add_rect(sq, 0, H - Inches(0.12), W, Inches(0.12), C_PRIMARY)
    add_txt(sq, Inches(0.5), Inches(1.8), Inches(12.3), Inches(1.5),
            'Q&A', 72, bold=True, color=C_PRIMARY, align=PP_ALIGN.CENTER)
    add_txt(sq, Inches(0.5), Inches(3.5), Inches(12.3), Inches(0.6),
            'Questions & Discussion', 18, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_txt(sq, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4),
            f"{app_name} · hello@nexyza.com", 10, color=C_MUTED, align=PP_ALIGN.CENTER)

    # ── Thank You slide ────────────────────────────────────────────────────────
    sty = prs.slides.add_slide(blank)
    add_bg(sty)
    add_rect(sty, 0, 0, W, Inches(0.12), C_PRIMARY)
    add_rect(sty, 0, H - Inches(0.12), W, Inches(0.12), C_PRIMARY)
    add_txt(sty, Inches(0.5), Inches(2.0), Inches(12.3), Inches(1.4),
            'Thank You', 60, bold=True, color=C_TEXT, align=PP_ALIGN.CENTER)
    add_txt(sty, Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5),
            f"Report generated by {app_name}", 14, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_txt(sty, Inches(0.5), Inches(4.2), Inches(12.3), Inches(0.5),
            'nexyza.com', 13, color=C_ACCENT, align=PP_ALIGN.CENTER)

    buf = _io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
